import streamlit as st
import pandas as pd
from pptx import Presentation
from docx import Document
from io import BytesIO

st.set_page_config(page_title="Générateur de PowerPoint")

# Affichage du logo
st.image("templates/logo-aon.jpg", width=150)

def generate_ppt(template_path, excel_file, sheet_name, client_selection, placeholders):
    """ Génère des présentations PowerPoint basées sur le modèle et les données """
    data = pd.read_excel(excel_file, sheet_name=sheet_name)
    filtered_data = data[data['client'] == client_selection]

    ppt_files = []

    for index, row in filtered_data.iterrows():
        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for placeholder, column in placeholders.items():
                                run.text = run.text.replace(placeholder, str(row[column]))

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        ppt_files.append((f"PROJET_REMISE_OFFRES_CONVENTIONS_{index + 1}.pptx", ppt_io))

    return ppt_files

def remplir_document_word(template_path, excel_file, sheet_name, client_selection, placeholders):
    # Lire les données Excel
    data = pd.read_excel(excel_file, sheet_name=sheet_name)
    filtered_data = data[data['client'] == client_selection]

    # Ouvrir le modèle Word
    doc = Document(template_path)

    # Parcourir les lignes du fichier Excel
    for index, row in filtered_data.iterrows():
        # Ajouter un saut de page pour chaque entrée
        doc.add_page_break()

        # Remplacer les espaces réservés dans le document Word
        for paragraph in doc.paragraphs:
            for key, value in placeholders.items():
                if f"{{{key}}}" in paragraph.text:
                    paragraph.text = paragraph.text.replace(f"{{{key}}}", str(row[key]))

    # Créer un objet BytesIO pour enregistrer le document
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)  # Rembobiner au début du flux

    return doc_io

# Interface utilisateur Streamlit
st.title("Générateur de PowerPoint")

# Chemins vers les modèles PowerPoint et Word
chemin_template_flottes = "templates/ppt_flottes.pptx"
chemin_template_mission = "templates/ppt_missions.pptx"
chemin_template_word = "templates/word.docx"

# Dictionnaires des espaces réservés pour les contrats "flottes" et "missions"
placeholders_flottes = {
    "client": "client", "date": "date", "nom": "nom", "adresse": "adresse",
    "cp": "cp", "ville": "ville", "leger": "leger", "lourd": "lourd",
    "semirem": "semirem", "reminf": "reminf", "deuxroues": "droue",
    "engins": "engins", "assureur": "ass", "echeann": "echeann",
    "frac": "frac", "reg": "reg", "fin": "fin"
}

placeholders_missions = {
    "client": "client", "date": "date", "nom": "nom", "adresse": "adresse",
    "cp": "cp", "ville": "ville", "assureur": "ass", "echeann": "echeann",
    "frac": "frac", "fin": "fin"
}

placeholders_word = {
    "client": "client", "date": "date", "adresse": "adresse",
    "cp": "cp", "ville": "ville", "assureur": "assureur", "camionette": "camionette",
    "camion": "camion", "deuxroues": "deuxroues", "engins": "engins", "autre": "autre",
    "effet": "effet", "siret": "siret", "activite": "activite", "risque": "risque"
}

# Téléchargement du fichier Excel
excel_file = st.file_uploader("Choisissez le fichier Excel", type=["xlsx", "xls"])

if excel_file is not None:
    # Lire les noms des feuilles du fichier Excel
    xls = pd.ExcelFile(excel_file)
    sheet_names = xls.sheet_names

    # Sélection de la feuille par l'utilisateur
    selected_sheet = st.selectbox("Choisissez l'onglet à traiter", sheet_names)

    # Lire les données de la feuille sélectionnée
    df = pd.read_excel(excel_file, sheet_name=selected_sheet)

    # Afficher les données de la feuille sélectionnée
    st.write(f"Données de l'onglet {selected_sheet}:")
    st.dataframe(df)

    # Sélection du client par l'utilisateur
    clients = df['client'].unique()
    client_selection = st.selectbox("Choisissez un client", clients)

    # Bouton pour générer les présentations "Flottes"
    if st.button("Générer PowerPoint Flottes"):
        ppt_files = generate_ppt(chemin_template_flottes, excel_file, selected_sheet, client_selection, placeholders_flottes)
        for filename, ppt_io in ppt_files:
            st.download_button(
                label=f"📥 Télécharger {filename}",
                data=ppt_io,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        st.balloons()  # Affiche une animation de ballons après la génération des présentations

    # Bouton pour générer les présentations "Missions"
    if st.button("Générer PowerPoint Mission"):
        ppt_files = generate_ppt(chemin_template_mission, excel_file, selected_sheet, client_selection, placeholders_missions)
        for filename, ppt_io in ppt_files:
            st.download_button(
                label=f"📥 Télécharger {filename}",
                data=ppt_io,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        st.balloons()  # Affiche une animation de ballons après la génération des présentations

    # Bouton pour générer le document Word
    if st.button("Générer Word"):
        word_file = remplir_document_word(chemin_template_word, excel_file, selected_sheet, client_selection, placeholders_word)
        st.download_button(
            label="📥 Télécharger le document Word",
            data=word_file,
            file_name="document_rempli.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        st.balloons()  # Affiche une animation de ballons après la génération du document Word
else:
    st.warning("Veuillez télécharger un fichier Excel.")
