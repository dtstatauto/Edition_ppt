import streamlit as st
import pandas as pd
from pptx import Presentation
from docx import Document
from io import BytesIO



def generate_ppt(template_path, excel_file, sheet_name, client_selection, placeholders):
    """ Génère des présentations PowerPoint basées sur le modèle et les données """
    data = pd.read_excel(excel_file, sheet_name=sheet_name)
    filtered_data = data[data['client'] == client_selection]

    ppt_files = []

    for index, row in filtered_data.iterrows():
        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for placeholder, column in placeholders.items():
                                run.text = run.text.replace(placeholder, str(row[column]))

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        ppt_files.append((f"PROJET_REMISE_OFFRES_CONVENTIONS_{index + 1}.pptx", ppt_io))

    return ppt_files

def remplir_document_word(template_path, excel_file, sheet_name, client_selection, placeholders):
    # Lire les données Excel
    data = pd.read_excel(excel_file, sheet_name=sheet_name)
    filtered_data = data[data['client'] == client_selection]

    # Ouvrir le modèle Word
    doc = Document(template_path)

    # Parcourir les lignes du fichier Excel
    for index, row in filtered_data.iterrows():
        # Ajouter un saut de page pour chaque entrée
        doc.add_page_break()

        # Remplacer les espaces réservés dans le document Word
        for paragraph in doc.paragraphs:
            for key, value in placeholders.items():
                if f"{{{key}}}" in paragraph.text:
                    paragraph.text = paragraph.text.replace(f"{{{key}}}", str(row[key]))

    # Créer un objet BytesIO pour enregistrer le document
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)  # Rembobiner au début du flux

    return doc_io